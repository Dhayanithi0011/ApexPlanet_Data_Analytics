"""
Task 4 — Presentation Generator
Creates a 10-slide PowerPoint presentation for the business audience.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================
# COLOR SCHEME (Dark Professional Theme)
# ============================================
BG_DARK      = RGBColor(0x0E, 0x11, 0x17)
BG_CARD      = RGBColor(0x16, 0x1B, 0x22)
ACCENT_CYAN  = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_PINK  = RGBColor(0xFF, 0x00, 0x7F)
ACCENT_GREEN = RGBColor(0x00, 0xFF, 0x66)
ACCENT_GOLD  = RGBColor(0xFF, 0xD7, 0x00)
TEXT_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY     = RGBColor(0xCC, 0xCC, 0xCC)
TEXT_LIGHT    = RGBColor(0xAA, 0xAA, 0xAA)

# ============================================
# HELPER FUNCTIONS
# ============================================
def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=TEXT_WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox

def add_shape_box(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

# ============================================
# CREATE PRESENTATION
# ============================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================
# SLIDE 1: TITLE
# ============================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, BG_DARK)

add_text_box(slide1, 1.5, 1.5, 10, 1.2,
             "Data Storytelling & Statistical Validation",
             font_size=36, color=ACCENT_CYAN, bold=True)

add_text_box(slide1, 1.5, 3.0, 10, 0.8,
             "Sales Performance and Customer Insights",
             font_size=22, color=TEXT_WHITE)

add_text_box(slide1, 1.5, 4.2, 10, 0.6,
             "ApexPlanet Data Analytics Internship | Task 4",
             font_size=16, color=TEXT_GRAY)

add_text_box(slide1, 1.5, 5.5, 10, 0.5,
             "Connecting What We Observed to What We Can Statistically Support",
             font_size=14, color=TEXT_LIGHT)

# ============================================
# SLIDE 2: BUSINESS OBJECTIVE
# ============================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)

add_text_box(slide2, 0.8, 0.4, 11, 0.8,
             "Business Objective",
             font_size=30, color=ACCENT_CYAN, bold=True)

add_shape_box(slide2, 0.8, 1.5, 11.5, 2.0, BG_CARD)
add_multi_text(slide2, 1.1, 1.7, 11, 1.8, [
    ("What We Set Out to Do", 20, ACCENT_GOLD, True),
    ("Analyze sales performance, identify important business patterns,", 16, TEXT_WHITE, False),
    ("and statistically validate one key finding.", 16, TEXT_WHITE, False),
])

add_shape_box(slide2, 0.8, 3.8, 11.5, 2.5, BG_CARD)
add_multi_text(slide2, 1.1, 4.0, 11, 2.3, [
    ("Our Approach", 20, ACCENT_GOLD, True),
    ("Task 1: Data Preparation and Cleaning", 15, TEXT_WHITE, False),
    ("Task 2: Exploratory Data Analysis and Visualization", 15, TEXT_WHITE, False),
    ("Task 3: Interactive Dashboard Development", 15, TEXT_WHITE, False),
    ("Task 4: Data Storytelling + Statistical Validation", 15, ACCENT_CYAN, True),
])

# ============================================
# SLIDE 3: WHAT WE LEARNED (Task 2/3 Findings)
# ============================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)

add_text_box(slide3, 0.8, 0.4, 11, 0.8,
             "What We Learned From the Analysis",
             font_size=30, color=ACCENT_CYAN, bold=True)

# KPI Cards
kpis = [
    ("1,000", "Transactions", ACCENT_CYAN),
    ("Rs.139.4M", "Total Revenue", ACCENT_GREEN),
    ("Rs.139,399", "Avg Transaction", ACCENT_GOLD),
    ("5", "Categories", ACCENT_PINK),
    ("8", "Cities", RGBColor(0x7B, 0x2C, 0xBF)),
]

for i, (val, label, color) in enumerate(kpis):
    x = 0.8 + i * 2.45
    add_shape_box(slide3, x, 1.5, 2.2, 1.6, BG_CARD)
    add_text_box(slide3, x + 0.1, 1.7, 2.0, 0.7, val, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide3, x + 0.1, 2.4, 2.0, 0.5, label, font_size=13, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_shape_box(slide3, 0.8, 3.5, 11.5, 3.2, BG_CARD)
add_multi_text(slide3, 1.1, 3.7, 11, 3.0, [
    ("Key Descriptive Findings", 20, ACCENT_GOLD, True),
    ("", 6, TEXT_WHITE, False),
    ("Electronics dominates total revenue at Rs.50.78M (36.4% of all revenue)", 15, TEXT_WHITE, False),
    ("with 354 transactions (35.4% of all orders).", 15, TEXT_WHITE, False),
    ("", 6, TEXT_WHITE, False),
    ("Average customer age: 41.4 years | Gender split: 51.1% Male, 48.9% Female", 15, TEXT_WHITE, False),
    ("", 6, TEXT_WHITE, False),
    ("Top city: Kolkata (Rs.19.52M) | Peak month: March (Rs.13.06M)", 15, TEXT_WHITE, False),
    ("", 6, TEXT_WHITE, False),
    ("Strong positive correlations: Quantity-Sales (r=0.65), Price-Sales (r=0.69)", 15, TEXT_WHITE, False),
])

# ============================================
# SLIDE 4: REVENUE STORY
# ============================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_DARK)

add_text_box(slide4, 0.8, 0.4, 11, 0.8,
             "The Revenue Story",
             font_size=30, color=ACCENT_CYAN, bold=True)

# Embed the revenue vs volume chart
if os.path.exists("revenue_vs_volume.png"):
    slide4.shapes.add_picture("revenue_vs_volume.png", Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))

add_text_box(slide4, 0.8, 1.4, 6, 0.5,
             "Electronics leads total revenue AND transaction count",
             font_size=14, color=ACCENT_GREEN, bold=True)

# ============================================
# SLIDE 5: AVERAGE TRANSACTION VALUE
# ============================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, BG_DARK)

add_text_box(slide5, 0.8, 0.4, 11, 0.8,
             "Average Transaction Value by Category",
             font_size=30, color=ACCENT_CYAN, bold=True)

if os.path.exists("hypothesis_test.png"):
    slide5.shapes.add_picture("hypothesis_test.png", Inches(0.8), Inches(1.4), Inches(7.0), Inches(5.0))

# Key question box
add_shape_box(slide5, 8.2, 1.4, 4.3, 5.0, BG_CARD)
add_multi_text(slide5, 8.5, 1.6, 3.8, 4.8, [
    ("The Question", 18, ACCENT_GOLD, True),
    ("", 6, TEXT_WHITE, False),
    ("Observed averages:", 14, TEXT_GRAY, False),
    ("Grocery: Rs.145,305", 14, TEXT_WHITE, False),
    ("Electronics: Rs.143,442", 14, TEXT_WHITE, False),
    ("Education: Rs.140,627", 14, TEXT_WHITE, False),
    ("Furniture: Rs.135,356", 14, TEXT_WHITE, False),
    ("Fashion: Rs.127,153", 14, TEXT_WHITE, False),
    ("", 6, TEXT_WHITE, False),
    ("Range: Rs.18,152", 14, ACCENT_PINK, True),
    ("", 6, TEXT_WHITE, False),
    ("Are these differences", 16, ACCENT_CYAN, True),
    ("statistically significant,", 16, ACCENT_CYAN, True),
    ("or just random noise?", 16, ACCENT_CYAN, True),
])

# ============================================
# SLIDE 6: HYPOTHESIS
# ============================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, BG_DARK)

add_text_box(slide6, 0.8, 0.4, 11, 0.8,
             "Business Hypothesis",
             font_size=30, color=ACCENT_CYAN, bold=True)

add_shape_box(slide6, 0.8, 1.5, 11.5, 1.2, BG_CARD)
add_multi_text(slide6, 1.1, 1.65, 11, 1.0, [
    ("Business Question", 18, ACCENT_GOLD, True),
    ("Do average sales values differ significantly across product categories?", 16, TEXT_WHITE, False),
])

add_shape_box(slide6, 0.8, 3.0, 5.5, 2.5, BG_CARD)
add_multi_text(slide6, 1.1, 3.2, 5.2, 2.3, [
    ("Null Hypothesis (H0)", 18, ACCENT_GREEN, True),
    ("", 4, TEXT_WHITE, False),
    ("Average sales are equal across", 15, TEXT_WHITE, False),
    ("all product categories.", 15, TEXT_WHITE, False),
    ("", 4, TEXT_WHITE, False),
    ("u_Education = u_Electronics = u_Fashion", 13, TEXT_GRAY, False),
    ("= u_Furniture = u_Grocery", 13, TEXT_GRAY, False),
])

add_shape_box(slide6, 6.8, 3.0, 5.5, 2.5, BG_CARD)
add_multi_text(slide6, 7.1, 3.2, 5.2, 2.3, [
    ("Alternative Hypothesis (H1)", 18, ACCENT_PINK, True),
    ("", 4, TEXT_WHITE, False),
    ("At least one category has a", 15, TEXT_WHITE, False),
    ("different average sales value.", 15, TEXT_WHITE, False),
    ("", 4, TEXT_WHITE, False),
    ("Significance Level: a = 0.05", 14, ACCENT_GOLD, True),
])

add_shape_box(slide6, 0.8, 5.8, 11.5, 1.2, BG_CARD)
add_multi_text(slide6, 1.1, 5.95, 11, 1.0, [
    ("Test Selection: Welch's ANOVA", 18, ACCENT_CYAN, True),
    ("Suitable for comparing means across multiple groups without assuming equal variances.", 15, TEXT_WHITE, False),
])

# ============================================
# SLIDE 7: STATISTICAL METHOD
# ============================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, BG_DARK)

add_text_box(slide7, 0.8, 0.4, 11, 0.8,
             "Statistical Method",
             font_size=30, color=ACCENT_CYAN, bold=True)

steps = [
    ("1. Data", "1,000 transactions across 5 categories"),
    ("2. Groups", "5 independent category groups"),
    ("3. Test", "Welch's ANOVA (unequal variance)"),
    ("4. Statistic", "F-statistic calculation"),
    ("5. P-value", "Probability of observed difference"),
    ("6. Decision", "Compare p-value to a = 0.05"),
]

for i, (step, desc) in enumerate(steps):
    y = 1.5 + i * 0.9
    add_shape_box(slide7, 0.8, y, 11.5, 0.75, BG_CARD)
    add_text_box(slide7, 1.1, y + 0.1, 2.5, 0.5, step, font_size=16, color=ACCENT_CYAN, bold=True)
    add_text_box(slide7, 3.8, y + 0.1, 8.2, 0.5, desc, font_size=15, color=TEXT_WHITE)

add_shape_box(slide7, 0.8, 7.0, 11.5, 0.35, RGBColor(0x1A, 0x23, 0x2E))

# ============================================
# SLIDE 8: RESULTS
# ============================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, BG_DARK)

add_text_box(slide8, 0.8, 0.4, 11, 0.8,
             "Results",
             font_size=30, color=ACCENT_CYAN, bold=True)

# Levene's test
add_shape_box(slide8, 0.8, 1.5, 3.5, 2.0, BG_CARD)
add_multi_text(slide8, 1.0, 1.65, 3.2, 1.8, [
    ("Levene's Test", 16, ACCENT_GOLD, True),
    ("", 4, TEXT_WHITE, False),
    ("Statistic: 0.9787", 14, TEXT_WHITE, False),
    ("P-value: 0.4182", 14, TEXT_WHITE, False),
    ("Variances: Equal", 14, ACCENT_GREEN, True),
])

# Main results
add_shape_box(slide8, 4.7, 1.5, 3.8, 2.0, BG_CARD)
add_multi_text(slide8, 4.9, 1.65, 3.5, 1.8, [
    ("Welch ANOVA", 16, ACCENT_GOLD, True),
    ("", 4, TEXT_WHITE, False),
    ("F-statistic: 0.7842", 14, TEXT_WHITE, False),
    ("P-value: 0.5359", 14, TEXT_WHITE, False),
    ("DF: (4.0, 429.6)", 14, TEXT_WHITE, False),
])

# Decision box
add_shape_box(slide8, 8.9, 1.5, 3.6, 2.0, BG_CARD)
add_multi_text(slide8, 9.1, 1.65, 3.3, 1.8, [
    ("Decision", 16, ACCENT_GOLD, True),
    ("", 4, TEXT_WHITE, False),
    ("0.5359 > 0.05", 14, TEXT_WHITE, False),
    ("", 4, TEXT_WHITE, False),
    ("FAIL TO REJECT H0", 15, ACCENT_PINK, True),
])

# Interpretation
add_shape_box(slide8, 0.8, 4.0, 11.5, 3.2, BG_CARD)
add_multi_text(slide8, 1.1, 4.2, 11, 3.0, [
    ("Interpretation", 20, ACCENT_GOLD, True),
    ("", 6, TEXT_WHITE, False),
    ("The p-value of 0.5359 is substantially greater than our significance level of 0.05.", 15, TEXT_WHITE, False),
    ("", 4, TEXT_WHITE, False),
    ("This means there is a 53.6% probability that we would observe differences this", 15, TEXT_WHITE, False),
    ("large (or larger) even if all category means were truly equal.", 15, TEXT_WHITE, False),
    ("", 4, TEXT_WHITE, False),
    ("Therefore, we do NOT have sufficient statistical evidence to conclude that", 15, ACCENT_CYAN, True),
    ("average transaction values differ across the five product categories.", 15, ACCENT_CYAN, True),
    ("", 4, TEXT_WHITE, False),
    ("The observed differences (Rs.127K to Rs.145K) could plausibly be due to", 15, TEXT_WHITE, False),
    ("random variation in the sample.", 15, TEXT_WHITE, False),
])

# ============================================
# SLIDE 9: BUSINESS MEANING
# ============================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9, BG_DARK)

add_text_box(slide9, 0.8, 0.4, 11, 0.8,
             "What This Means for Business",
             font_size=30, color=ACCENT_CYAN, bold=True)

# The Key Distinction
add_shape_box(slide9, 0.8, 1.5, 11.5, 2.0, BG_CARD)
add_multi_text(slide9, 1.1, 1.65, 11, 1.8, [
    ("The Key Distinction: Total Revenue vs Average Transaction Value", 18, ACCENT_GOLD, True),
    ("", 4, TEXT_WHITE, False),
    ("Electronics is #1 in Total Revenue (Rs.50.78M) and #1 in Transaction Volume (354 orders).", 15, TEXT_WHITE, False),
    ("However, its average transaction value (Rs.143,442) is NOT statistically different from other categories.", 15, ACCENT_CYAN, True),
])

# Comparison table
add_shape_box(slide9, 0.8, 3.8, 5.5, 3.2, BG_CARD)
add_multi_text(slide9, 1.1, 3.95, 5.2, 3.0, [
    ("Electronics Revenue Leadership", 17, ACCENT_GREEN, True),
    ("", 4, TEXT_WHITE, False),
    ("Total Revenue:  Rs.50.78M (#1)", 14, TEXT_WHITE, False),
    ("Transactions:     354 (#1, 35.4% of orders)", 14, TEXT_WHITE, False),
    ("Avg Value:        Rs.143,442 (#2, not significant)", 14, TEXT_WHITE, False),
    ("", 4, TEXT_WHITE, False),
    ("Revenue leadership is DRIVEN BY VOLUME,", 14, ACCENT_GOLD, True),
    ("not by statistically higher transaction values.", 14, ACCENT_GOLD, True),
])

add_shape_box(slide9, 6.8, 3.8, 5.5, 3.2, BG_CARD)
add_multi_text(slide9, 7.1, 3.95, 5.2, 3.0, [
    ("Practical Implication", 17, ACCENT_PINK, True),
    ("", 4, TEXT_WHITE, False),
    ("Management should NOT assume that a", 14, TEXT_WHITE, False),
    ("category has inherently higher transaction", 14, TEXT_WHITE, False),
    ("value based on observed sample means.", 14, TEXT_WHITE, False),
    ("", 4, TEXT_WHITE, False),
    ("The business should focus on:", 14, ACCENT_CYAN, True),
    ("Maintaining transaction volume", 14, TEXT_WHITE, False),
    ("Understanding volume drivers", 14, TEXT_WHITE, False),
    ("Collecting more data for validation", 14, TEXT_WHITE, False),
])

# ============================================
# SLIDE 10: RECOMMENDATIONS
# ============================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10, BG_DARK)

add_text_box(slide10, 0.8, 0.4, 11, 0.8,
             "Recommendations",
             font_size=30, color=ACCENT_CYAN, bold=True)

recs = [
    ("1", "Protect Electronics Volume",
     "Maintain and grow the transaction count that drives revenue leadership.", ACCENT_GREEN),
    ("2", "Investigate Volume Drivers",
     "Understand why Electronics attracts significantly more transactions than other categories.", ACCENT_CYAN),
    ("3", "Avoid Over-Interpreting Averages",
     "Observed category mean differences were not statistically significant. Do not make pricing or inventory decisions based solely on these observed averages.", ACCENT_PINK),
    ("4", "Collect More Data",
     "A larger dataset over a longer period could provide stronger statistical evidence. The current sample of 1,000 transactions may be insufficient.", ACCENT_GOLD),
    ("5", "Continue Statistical Validation",
     "Use hypothesis testing before making major business decisions based on observed data patterns.", RGBColor(0x7B, 0x2C, 0xBF)),
]

for i, (num, title, desc, color) in enumerate(recs):
    y = 1.4 + i * 1.15
    add_shape_box(slide10, 0.8, y, 11.5, 1.0, BG_CARD)
    add_text_box(slide10, 1.0, y + 0.1, 0.5, 0.5, num, font_size=20, color=color, bold=True)
    add_text_box(slide10, 1.6, y + 0.05, 3.0, 0.4, title, font_size=15, color=TEXT_WHITE, bold=True)
    add_text_box(slide10, 1.6, y + 0.45, 10.5, 0.5, desc, font_size=13, color=TEXT_GRAY)

# ============================================
# SAVE
# ============================================
prs.save("presentation.pptx")
print("Presentation saved -> presentation.pptx")
